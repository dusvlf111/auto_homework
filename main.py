import fitz  # PyMuPDF               #pdf
# import hwp5py                       #hwp
from pptx import Presentation  #ppt
import docx  #docx
import os
import json
import openai
#import pprint

# config.json 파일을 찾고 설정하는 코드
# config.json 파일잉 없다면 write_config_file


class Config:

  def __init__(self):
    self.CONFIG_FILE_PATH = os.path.join(os.path.dirname(__file__),
                                         "config.json")

  def read_config_file(self):
    if not os.path.exists(self.CONFIG_FILE_PATH):
      self.write_config_file({"file_path": "", "api_key": ""})
    with open(self.CONFIG_FILE_PATH, "r") as f:
      config = json.load(f)
    return config

  def write_config_file(self, config):
    with open(self.CONFIG_FILE_PATH, "w") as f:
      json.dump(config, f, indent=4)

  def get_folder_path(self):
    config = self.read_config_file()
    file_path = config.get("file_path", "")
    if not os.path.isfile(file_path):
      user_input = input(
        "Enter the file path (or enter 1 to save to the same directory as the main file): "
      )
      if user_input == '1':
        main_file_path = os.path.abspath(__file__)
        directory = os.path.dirname(main_file_path)
        file_path = os.path.join(directory, "input_file")
      else:
        file_path = user_input
      config["file_path"] = file_path
      self.write_config_file(config)
    return file_path
  
  def get_file_path(self):
    config = self.read_config_file()
    file_path = config.get("file_path", "")
    return file_path
    


  def get_api_key(self):
    config = self.read_config_file()
    api_key = config.get("api_key", "")
    if not api_key:
      api_key = input("Enter the API key: ")
      config["api_key"] = api_key
      self.write_config_file(config)
    return api_key

  def run(self):
    file_path = self.get_folder_path()
    api_key = self.get_api_key()
    print(f"Configured file path: {file_path}")
    print(f"Configured API key: {api_key}")
    


class ChangeText:

  def __init__(self):
    self.config = Config()
    self.folder_path = self.config.get_file_path()

  def print_files_in_folder(self):
    file_list = []
    for idx, (root, dirs, files) in enumerate(os.walk(self.folder_path)):
      for file in files:
        file_path = os.path.join(root, file)
        file_extension = os.path.splitext(file)[1]

        if os.path.isfile(file_path):
          file_list.append(file_path)
          print(
            f"{len(file_list)}. File path: {file_path.replace('//' ,'/')}, File extension: {file_extension}"
          )

    while True:
      try:
        selected_file = int(
          input("Enter the number of the file you want to select: "))
        if selected_file == 0:
          return None
        self.file = file_list[selected_file - 1]
        #리스트에서 선택한 인덱스값의 파일경로 출력,
        return self.file
    
      except (ValueError, IndexError):
        print("Invalid input, please try again.")

  def extract_pdf_text(self):
    with fitz.open(self.file) as doc:
      text = ""
      for page in doc:
        text += page.get_text()
      return text

  # def extract_hwp_text(self):
  #     with hwp5py.HWP5File(self.file) as hwp:
  #         return hwp.to_text()

  def extract_ppt_text(self):
    prs = Presentation(self.file)
    text_pages = ""
    for slide in prs.slides:
        text = ""
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text += shape.text+"\n"
        text_pages+=(text+"\n")
    return text_pages


  def extract_word_text(self):
    doc = docx.Document(self.file)
    text = []
    for para in doc.paragraphs:
      text.append(para.text)
    return "\n".join(text)

  def extract_text(self):
    with open(self.file, "r") as f:
      text = f.read()
    return text

  def run(self):
    
    self.file = self.print_files_in_folder()
    extension = os.path.splitext(self.file)[1]

    if extension == ".pdf":
      text = self.extract_pdf_text()
    elif extension == ".hwp":
      text = self.extract_hwp_text()
    elif extension in [".ppt", ".pptx"]:
      text = self.extract_ppt_text()
    elif extension == ".docx":
      text = self.extract_word_text()
    elif extension == ".txt":
      text = self.extract_text()
    else:
      raise ValueError("Unsupported file type")

    return text


class GPT:

  def __init__(self, prompt_list):
    self.prompt_list = prompt_list
    self.config = Config()
    openai.api_key = self.config.get_api_key()

  def chat_with_gpt1(self):
    prompt_sentences = self.prompt_list.split('\n')
     # 입력된 prompt를 문장 단위로 쪼갬
    response = ""
    print("-"*90,"\n"*3)
    for sentence in prompt_sentences:
        sentence.strip()
        if len(sentence) <= 10 : continue
        response = openai.Completion.create(
            engine="text-davinci-003",
            prompt=sentence,
            max_tokens=1000,
            n=1,
            stop=None,
            temperature=0.5
        ).choices[0].text
        print("질문: \n \n",sentence)
        print("-"*90)
        print("지피티 답변:\n",response)
        print("-"*90)
        
  def chat_with_gpt(self):
    print("write option:")
    self.option = str(input().strip())
    response = openai.Completion.create(
      engine="text-davinci-003",
      prompt=str(self.prompt_list)+"\n details:  "+self.option,
      max_tokens=2048,
      n=1,
      stop=None,
      temperature=0.5,
    )
    print("-"*90)
    print(self.prompt_list,"\n details:  "+self.option )
    print("-"*90)
    print("\n 지피티의 답변: \n",response.choices[0].text)



if __name__ == "__main__":
  config = Config()
  change = ChangeText()
  a = input("바로 입력하려면 1 아니면 엔터:\n ")
  if a == "1":
    print("입력하시오:")
    GPT(input().strip()).chat_with_gpt()
  else:  
    text = change.run()
    print(text)
    try:
      gpt_obj = GPT(text)
      gpt_obj.chat_with_gpt()
    except:
      gpt_obj = GPT(text)
      gpt_obj.chat_with_gpt1()
 

